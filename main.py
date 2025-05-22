import os
from flask import Flask, render_template, request, session, send_file
from openai import OpenAI
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from io import BytesIO

app = Flask(__name__)
app.secret_key = os.getenv("FLASK_SECRET_KEY", "super-secret-dev-key")

client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

STYLE_PROMPTS = {
    '3-paragraph Narrative (Default)':
    """
You are a master sales strategist at BiteSpeed.

The merchant you're pitching has the following profile:
Industry: {industry}
AOV: {aov}
Geography: {geo}
Store Maturity: {maturity}
Primary Goal: {goal}

Write a 3-paragraph, high-impact sales pitch that:
1) Hooks by naming their industry and region and paints the risk of inaction.
2) Explains the top 2–3 BiteSpeed features tied to their goal + AOV, and drops in one real BiteSpeed case study from that same industry & region with a hard metric.
3) Recommends the precise plan (Free/Pro/Omnichannel), creates urgency with a time-bound offer, and closes with “Ready to start your free trial?”

Do not use emojis, bullets, or headings—just three polished paragraphs.
""",
    '5-Line WhatsApp Pitch':
    """
You are writing a 5-line WhatsApp pitch for a merchant.
Their profile is:
Industry: {industry}
AOV: {aov}
Geography: {geo}
Store Maturity: {maturity}
Primary Goal: {goal}

Make sure the pitch clearly focuses on how BiteSpeed addresses their primary goal.
Each line should be impactful, relevant to their industry and goal, and encourage immediate action.
Conclude with: “Want me to set this up for you?”

Do not use emojis.
""",
    'Cold Email Template':
    """
You're a growth marketer at BiteSpeed. Your target merchant has the following profile:
Industry: {industry}
AOV: {aov}
Geography: {geo}
Store Maturity: {maturity}
Primary Goal: {goal}

Write a cold email that includes:
- A catchy subject line
- A 2-line hook naming their industry and goal
- A 3-line value proposition linking their AOV/goal to specific BiteSpeed features
- A short client case study (metric + brand)
- A call to action to start a free trial or book a demo

Do not use emojis. Tone: Friendly and confident, not overly pushy.
""",
    'Deck Slide Suggestions':
    """
You are preparing pitch deck slides for a merchant with the following profile:
Industry: {industry}
AOV: {aov}
Geography: {geo}
Store Maturity: {maturity}
Primary Goal: {goal}

Suggest 4 slide titles and for each:
- Provide 2 persuasive bullet points relevant to their profile
- Use benefit-driven language
- Include one metric or proof point in at least one slide

Do not use emojis.
""",
    'Formal vs Casual Comparison':
    """
You are writing 2 versions of the same pitch for a merchant with this profile:
Industry: {industry}
AOV: {aov}
Geography: {geo}
Store Maturity: {maturity}
Primary Goal: {goal}

Write two versions of a 5-line pitch:
1) Formal and professional (for B2B email)
2) Casual and punchy (for WhatsApp)

Both should:
- Start with a hook tied to the merchant’s primary goal
- Mention a key BiteSpeed feature or outcome based on their profile
- End with a call to action

Do not use emojis.
"""
}

def set_slide_title(slide, text):
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = text
        return
    for shape in slide.placeholders:
        if shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.TITLE:
            shape.text = text
            return


def find_body_placeholder(slide):
    for shape in slide.placeholders:
        if shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.BODY:
            return shape
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape is not slide.shapes.title:
            return shape
    return None


@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'GET':
        session.pop('messages', None)
        session.pop('industry', None)
        session.pop('goal', None)

    session.setdefault('messages', [])

    if request.method == 'POST':
        user_refine = request.form.get('user_message', '').strip()
        if user_refine and session['messages']:
            msgs = session['messages']
            msgs.append({'role': 'user', 'content': user_refine})
            resp = client.chat.completions.create(model="gpt-4o-mini",
                                                  messages=[{
                                                      'role':
                                                      m['role'],
                                                      'content':
                                                      m['content']
                                                  } for m in msgs],
                                                  temperature=0.7)
            reply = resp.choices[0].message.content.strip()
            msgs.append({'role': 'assistant', 'content': reply})
            session['messages'] = msgs
        else:
            data = {
                'industry': request.form['industry'],
                'aov': request.form['aov'],
                'geo': request.form['geo'],
                'maturity': request.form['maturity'],
                'goal': request.form['goal'],
                'style': request.form['style']
            }
            session['industry'] = data['industry']
            session['goal'] = data['goal']
            prompt_template = STYLE_PROMPTS.get(
                data['style'],
                STYLE_PROMPTS['3-paragraph Narrative (Default)'])
            prompt = prompt_template.format(**data)
            resp = client.chat.completions.create(model="gpt-4o-mini",
                                                  messages=[{
                                                      'role': 'user',
                                                      'content': prompt
                                                  }],
                                                  temperature=0.7)
            pitch = resp.choices[0].message.content.strip()
            session['messages'] = [{'role': 'assistant', 'content': pitch}]

    return render_template('index.html')


@app.route('/download_ppt', methods=['POST'])
def download_ppt():
    return "(Placeholder) PPT download feature to be implemented or re-enabled."


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 8080))
    app.run(host='0.0.0.0', port=port, debug=True)
